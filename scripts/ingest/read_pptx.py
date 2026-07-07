#!/usr/bin/env python3
"""Print a .pptx slide by slide (for citation as [file, slide N: "title"]).

Usage: python3 scripts/ingest/read_pptx.py <file.pptx>
Exists so the workshop allowlist can pre-approve THIS script instead of arbitrary python.
"""
import sys

try:
    from pptx import Presentation
except ImportError:
    sys.exit("Missing python-pptx. Install: pip install python-pptx  (or use the workshop venv)")

if len(sys.argv) != 2:
    sys.exit(__doc__)

prs = Presentation(sys.argv[1])
for n, slide in enumerate(prs.slides, 1):
    print(f"--- slide {n} ---")
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            print(shape.text_frame.text)
